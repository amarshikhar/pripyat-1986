#!/usr/bin/env python3
"""
Generate PRIPYAT-1986 Hackathon Presentation — WattAgents
Targets 85-90/100 on triage report by addressing every identified gap.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ═══════════════════════════════════════════════════════════════════════
# COLOR PALETTE — Capgemini-inspired
# ═══════════════════════════════════════════════════════════════════════
NAVY       = RGBColor(0x0C, 0x23, 0x40)
DARK_BLUE  = RGBColor(0x1B, 0x36, 0x5F)
TEAL       = RGBColor(0x00, 0x9B, 0x8D)
ACCENT     = RGBColor(0x00, 0x70, 0xAD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY      = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY      = RGBColor(0x33, 0x33, 0x33)
MGRAY      = RGBColor(0x77, 0x77, 0x77)
RED        = RGBColor(0xCC, 0x00, 0x00)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
GREEN      = RGBColor(0x22, 0x8B, 0x22)
YELLOW_D   = RGBColor(0xBB, 0x99, 0x00)
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xF4)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_RED  = RGBColor(0xFD, 0xE8, 0xE8)
LIGHT_ORANGE = RGBColor(0xFD, 0xF0, 0xE0)
LIGHT_GREEN  = RGBColor(0xE8, 0xF5, 0xE8)

SW = Inches(13.333)
SH = Inches(7.5)
MARGIN = Inches(0.5)
CW = SW - 2 * MARGIN  # content width

# ═══════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════
def rect(slide, l, t, w, h, color, border=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not border:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame

def para(tf, text, sz=10, color=DGRAY, bold=False, align=PP_ALIGN.LEFT,
         font='Calibri', space_before=Pt(1), space_after=Pt(1)):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.space_before = space_before; p.space_after = space_after
    return p

def bullet(tf, text, level=0, sz=10, color=DGRAY, bold=False, font='Calibri'):
    p = tf.add_paragraph()
    p.text = text; p.level = level; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font
    p.space_before = Pt(1); p.space_after = Pt(1)
    return p

def header_bar(slide, title, subtitle=None):
    bh = Inches(1.05) if not subtitle else Inches(1.25)
    bar = rect(slide, Inches(0), Inches(0), SW, bh, NAVY)
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = MARGIN; tf.margin_top = Inches(0.15)
    tf.margin_right = MARGIN
    para(tf, title, sz=22, color=WHITE, bold=True)
    if subtitle:
        para(tf, subtitle, sz=13, color=RGBColor(0x88, 0xBB, 0xDD))
    # Teal accent strip
    rect(slide, Inches(0), bh, SW, Pt(3), TEAL)
    return bh + Pt(3)

def footer(slide, num):
    fy = SH - Inches(0.32)
    rect(slide, MARGIN, fy - Inches(0.02), CW, Pt(0.75), RGBColor(0xDD, 0xDD, 0xDD))
    tf = txbox(slide, MARGIN, fy, Inches(7), Inches(0.28))
    para(tf, 'Company Confidential © Capgemini 2026. All rights reserved.', sz=7, color=MGRAY)
    tf2 = txbox(slide, SW - Inches(3), fy, Inches(2.5), Inches(0.28))
    para(tf2, f'WattAgents — PRIPYAT-1986 | {num}/20', sz=7, color=MGRAY, align=PP_ALIGN.RIGHT)

def add_table(slide, rows_data, left, top, width, col_widths=None, hdr_color=TEAL):
    """rows_data: list of lists. First row = headers."""
    nrows = len(rows_data); ncols = len(rows_data[0])
    row_h = Inches(0.3)
    tbl_h = row_h * nrows
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, tbl_h)
    tbl = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ''
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]; p.text = str(val)
            p.font.name = 'Calibri'
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hdr_color
                p.font.size = Pt(9); p.font.color.rgb = WHITE; p.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LGRAY if r % 2 == 0 else WHITE
                p.font.size = Pt(8); p.font.color.rgb = DGRAY
    return shape

def code_box(slide, text, left, top, width, height, bg=RGBColor(0xF8, 0xF8, 0xF8)):
    box = rrect(slide, left, top, width, height, bg)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.08)
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line; p.font.size = Pt(9); p.font.name = 'Consolas'
        p.font.color.rgb = DGRAY; p.space_before = Pt(0); p.space_after = Pt(0)
    return box

def kpi_box(slide, value, label, left, top, val_color=TEAL):
    w, h = Inches(2.8), Inches(0.85)
    box = rrect(slide, left, top, w, h, LIGHT_TEAL)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.05)
    para(tf, value, sz=18, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    para(tf, label, sz=8, color=DGRAY, align=PP_ALIGN.CENTER)
    return box

def section_box(slide, title, items, left, top, width, height, title_bg=TEAL, bg=WHITE):
    """Colored header box + content area"""
    hdr_h = Inches(0.3)
    hdr = rect(slide, left, top, width, hdr_h, title_bg)
    tf = hdr.text_frame; tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.02)
    para(tf, title, sz=10, color=WHITE, bold=True)
    content = rrect(slide, left, top + hdr_h, width, height - hdr_h, bg)
    tf2 = content.text_frame; tf2.word_wrap = True
    tf2.margin_left = Inches(0.1); tf2.margin_top = Inches(0.05)
    tf2.margin_right = Inches(0.05)
    for item in items:
        bullet(tf2, item, sz=9)
    return content

# ═══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title Slide — Full navy background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rect(slide, Inches(0), Inches(0), SW, SH, NAVY)
    # Teal accent stripe
    rect(slide, Inches(0), Inches(2.8), SW, Pt(4), TEAL)

    tf = txbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.6))
    para(tf, 'Microsoft Agentic Industry Hackathon 2026', sz=14, color=TEAL, bold=True,
         align=PP_ALIGN.CENTER)

    tf2 = txbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2))
    para(tf2, 'PRIPYAT-1986', sz=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    para(tf2, 'Agentic AI Crisis Response Architecture', sz=20, color=RGBColor(0xAA, 0xCC, 0xEE),
         align=PP_ALIGN.CENTER)

    # Track badge
    badge = rrect(slide, Inches(4.5), Inches(3.1), Inches(4.3), Inches(0.45), TEAL)
    btf = badge.text_frame
    para(btf, 'TRACK A: GitHub + Azure  |  Industry: Utilities Sector', sz=13, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)

    tf3 = txbox(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.5))
    para(tf3, 'Real-time Resilience Approach — Operations Assistant', sz=16, color=WHITE,
         align=PP_ALIGN.CENTER)

    tf4 = txbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5))
    para(tf4, 'Team WattAgents', sz=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    tf5 = txbox(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5))
    para(tf5, 'Shruti Kate  •  Amar Shikhar  •  V. Sreenivas  •  Deendayalan Balaji  •  M. Navenaa  •  Shyamdinesh',
         sz=11, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    tf6 = txbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4))
    para(tf6, 'April 2026', sz=11, color=MGRAY, align=PP_ALIGN.CENTER)


def slide_02_exec_summary(prs):
    """Executive Summary — Use Case"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Executive Summary', 'Use Case: Real-time Resilience — Utilities Sector  |  Track A: GitHub + Azure')

    # WHAT box
    bx = Inches(0.5); by = y + Inches(0.1)
    bw = Inches(6.1); bh = Inches(2.2)
    section_box(slide, 'WHAT — Problem & Solution', [
        'Multi-agent AI system replaying the 1986 Chernobyl disaster through 6 autonomous agents',
        'Monitors real-time reactor telemetry, detects compound safety violations across 6 dimensions',
        'Issues automated shutdown (AZ-5 SCRAM) and evacuation orders — decisions operators refused to make',
        'Proves agentic AI prevents catastrophic industrial failure in safety-critical infrastructure',
        'NOT a nuclear simulator — a reusable safety-critical AI architecture (Chernobyl = proof dataset)',
    ], bx, by, bw, bh, title_bg=ACCENT)

    # WHY box
    section_box(slide, 'WHY — The "Dyatlov Effect"', [
        'Human operators under cognitive pressure exhibit confirmation bias + authority gradient',
        'No existing system detects compound failure modes across simultaneous dimensions',
        'Chain-of-command paralysis delayed Chernobyl evacuation by 35+ hours',
        'AI eliminates cognitive bias from life-or-death decisions in real-time',
    ], bx + bw + Inches(0.15), by, Inches(6.1), bh, title_bg=DARK_BLUE)

    # WHO box
    by2 = by + bh + Inches(0.1)
    section_box(slide, 'WHO — Stakeholders', [
        'Plant operators & shift supervisors (primary users)',
        'IAEA / national nuclear regulatory authorities',
        'Civil emergency coordinators & evacuation planners',
        'Civilian populations in plant exclusion zones (49,000+ residents)',
        'Enterprise clients: utilities, oil & gas, aviation, pharma (reusable architecture)',
    ], bx, by2, Inches(6.1), Inches(1.75), title_bg=TEAL)

    # HOW box
    section_box(slide, 'HOW — 6-Agent Pipeline', [
        'Sensor → Risk → Decision → Dyatlov (adversarial) → Evacuation → Comms',
        'Guardrail-Union: Rules FIRST, AI SECOND, Final = UNION (AI adds, never removes)',
        '70% Azure OpenAI GPT-4o + 30% rule-based formula for risk scoring',
        'Hard guardrails non-negotiable — LLM advisory layer enhances safety floor',
        'Graceful degradation: works identically without API keys (100% rule-based)',
    ], bx + Inches(6.25), by2, Inches(6.1), Inches(1.75), title_bg=NAVY)

    # KPIs
    ky = by2 + Inches(1.85) + Inches(0.1)
    kpi_box(slide, '38 min Earlier', 'Emergency Shutdown (vs. explosion at 01:23:40)', Inches(0.5), ky)
    kpi_box(slide, '35 hrs Earlier', 'Civilian Evacuation (vs. Soviet order Apr 27 14:00)', Inches(3.5), ky)
    kpi_box(slide, '100%', 'Safety Protocol Compliance (vs. 0% in 1986)', Inches(6.5), ky)
    kpi_box(slide, '<$0.10/run', 'Cost Per Simulation on Azure OpenAI', Inches(9.5), ky)

    footer(slide, 2)


def slide_03_not_a_simulator(prs):
    """Reusability Thesis — The Architecture is the Product"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, '"Not a Nuclear Simulator" — Reusable Safety-Critical AI Architecture')

    # Big quote
    qbox = rrect(slide, Inches(0.5), y + Inches(0.1), CW, Inches(1.0), LIGHT_TEAL)
    qtf = qbox.text_frame; qtf.word_wrap = True
    qtf.margin_left = Inches(0.2); qtf.margin_top = Inches(0.05)
    para(qtf, '"PRIPYAT-1986 is not a nuclear simulator. It\'s a reusable safety-critical AI architecture that uses '
         'Chernobyl as its most dramatic training dataset. Connect it to your SCADA API, and it runs on your grid '
         'tomorrow. Connect it to your BSEE feed, and it monitors your wells. The architecture is the product. '
         'The disaster is just the proof."', sz=11, color=NAVY, bold=True, align=PP_ALIGN.LEFT)

    # Cross-industry table
    ty = y + Inches(1.25)
    rows = [
        ['PRIPYAT-1986 Agent', 'Oil & Gas', 'Aviation', 'Pharmaceuticals', 'Power Grid'],
        ['SensorAgent\n(threshold anomaly)', 'SCADA well pressure\nmonitor', 'Flight parameter\nenvelope checker',
         'Batch process\ntemp/pH monitor', 'PMU voltage/\nfrequency monitor'],
        ['RiskAgent\n(rule + LLM blend)', 'Well blowout\nprobability scorer', 'Bird strike / turbulence\nrisk scorer',
         'Contamination /\ncross-reaction risk', 'N-1 contingency\nanalyser'],
        ['DecisionAgent\n(guardrails + AI)', 'Emergency well\nshutdown', 'Go-around /\ndiversion decision',
         'Batch rejection /\nline lockdown', 'Load shedding /\nislanding decision'],
        ['EvacuationAgent\n(logistics)', 'Offshore rig\nevacuation planner', 'Airport terminal\nevacuation',
         'Clean room lockdown\n+ personnel routing', 'Rolling blackout\nzone planner'],
        ['CommsAgent\n(regulatory)', 'BSEE incident report\n+ coast guard', 'ATC alert +\nNTSB notification',
         'FDA deviation report\n+ recall notice', 'NERC event report\n+ ISO notification'],
        ['DyatlovAgent\n(adversarial pressure)', 'Production manager\noverride pressure', 'Captain authority\ngradient',
         'Manufacturing quota\npressure', 'Grid dispatcher\ncongestion pressure'],
    ]
    add_table(slide, rows, Inches(0.5), ty, CW)

    # Customization points
    cy = ty + Inches(2.5)
    tf = txbox(slide, Inches(0.5), cy, CW, Inches(0.3))
    para(tf, 'Customization Points — To Adapt for a New Industry:', sz=12, color=NAVY, bold=True)
    rows2 = [
        ['Module to Change', 'What Changes', 'What Stays (Reusable As-Is)'],
        ['timeline_data.py', 'Replace ReactorState → domain dataclass (WellState, FlightState, etc.)', 'Data flow pattern, pipeline structure'],
        ['config.py THRESHOLDS', 'Replace nuclear thresholds → domain values (PSI, altitude, pH)', 'Config structure, risk level definitions'],
        ['llm_client.py prompts', 'Replace nuclear expert persona → domain expert', 'JSON schema pattern, fallback logic, timeout handling'],
        ['agents.py checks/actions', 'Replace nuclear checks → domain checks; SCRAM → SHUTDOWN', 'Agent interface, message bus, guardrail-union pattern'],
        ['static/* dashboard', 'Replace chart labels and units', 'Dashboard architecture, WebSocket streaming, Plotly.js'],
    ]
    add_table(slide, rows2, Inches(0.5), cy + Inches(0.3), CW)
    # Reuse callout
    rbox = rrect(slide, Inches(0.5), cy + Inches(2.15), CW, Inches(0.55), LIGHT_BLUE)
    rtf = rbox.text_frame; rtf.word_wrap = True; rtf.margin_left = Inches(0.15)
    para(rtf, 'Everything else transfers as-is: Orchestrator, Message Bus, LLM Client, Guardrail-Union Engine, '
         'Dual Timeline Engine, Dashboard Framework, Evaluation Pipeline, IaC templates, Security model.',
         sz=10, color=NAVY, bold=True)
    footer(slide, 3)


def slide_04_use_case(prs):
    """Use Case Detail — Real-time Resilience"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Use Case: Real-time Resilience — Utilities Sector',
                   'Preventing Catastrophic Failure with Agentic AI Safety Co-Pilot')

    # Problem → Solution → Impact flow
    pw = Inches(4.0)
    px = Inches(0.5)
    py = y + Inches(0.1)
    ph = Inches(2.6)

    section_box(slide, 'PROBLEM', [
        'Operators suffer confirmation bias under pressure ("Dyatlov Effect")',
        'No compound-risk detection across simultaneous failure modes',
        'ECCS manually disabled, only 6-8 control rods at explosion (need 30)',
        'Xenon poisoning + rod withdrawal + ECCS bypass = invisible compound risk',
        '35+ hour evacuation delay due to chain-of-command paralysis',
        'Single threshold alarms miss converging multi-dimensional failures',
    ], px, py, pw, ph, title_bg=RED, bg=LIGHT_RED)

    section_box(slide, 'SOLUTION — 6-Agent AI Pipeline', [
        'SensorAgent monitors 6 dimensions against RBMK-1000 limits continuously',
        'RiskAgent blends LLM reasoning (70%) + deterministic rules (30%)',
        'DecisionAgent applies guardrail-union: rules FIRST, AI adds, never removes',
        'DyatlovAgent adversarially tests decisions (historically accurate pushback)',
        'EvacuationAgent plans logistics: 49,000 people, 1,200 buses, 3 routes',
        'CommsAgent broadcasts + Gaussian plume radiation modeling (5-300km)',
    ], px + pw + Inches(0.15), py, pw, ph, title_bg=ACCENT, bg=LIGHT_BLUE)

    section_box(slide, 'IMPACT — Measurable KPIs', [
        'AI issues AZ-5 SCRAM at 00:45 — 38 min before the actual explosion',
        'Evacuation ordered at 01:30 — 35 hrs before Soviet government',
        '100% safety protocol compliance (vs. 0% historically)',
        'Compound risk score detects invisible multi-dimensional failure',
        'Cost: <$0.10 per full simulation run on Azure OpenAI',
        'Architecture reusable across oil & gas, aviation, pharma, power grid',
    ], px + 2 * (pw + Inches(0.15)), py, pw, ph, title_bg=GREEN, bg=LIGHT_GREEN)

    # Alignment callout
    ay = py + ph + Inches(0.15)
    abox = rrect(slide, MARGIN, ay, CW, Inches(0.55), LIGHT_TEAL)
    atf = abox.text_frame; atf.word_wrap = True; atf.margin_left = Inches(0.15)
    para(atf, 'Hackathon Alignment: Maps directly to Utilities Sector — Real-time Resilience Approach (Operations '
         'Assistant). Control-room supervision, real-time risk detection, rapid response. Same pattern as MISO + '
         'Microsoft Azure AI grid platform deployed January 2026 for 45 million people across 15 states.',
         sz=10, color=NAVY, bold=True)

    # AI counterfactual decisions
    dy = ay + Inches(0.7)
    tf = txbox(slide, MARGIN, dy, CW, Inches(0.3))
    para(tf, 'AI Counterfactual Decisions — What Our System Would Have Done:', sz=11, color=NAVY, bold=True)
    rows = [
        ['Time', 'Historical (Human)', 'AI Decision', 'Lead Time'],
        ['00:28', 'Continued despite 30 MW + ECCS disabled', 'SCRAM — xenon poisoning detected', '55 min before explosion'],
        ['00:43', 'Increased power with 6-8 rods', 'REJECT power increase — rods below safe minimum', '41 min'],
        ['01:00', 'Proceeded with test, no conditions met', 'ABORT TEST — unsafe test conditions', '24 min'],
        ['01:19', 'Pressed AZ-5 at 01:23:40 (too late)', 'AZ-5 triggered when rods dropped below 15', '4 min 40s'],
        ['01:30', 'Delayed evacuation until Apr 27 14:00', 'IMMEDIATE EVACUATION — radiation confirms breach', '35+ hours'],
    ]
    add_table(slide, rows, MARGIN, dy + Inches(0.3), CW)
    footer(slide, 4)


def slide_05_timeline(prs):
    """Chernobyl Timeline & Dual Engine"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Chernobyl Timeline — 20 INSAG-7 Events & Dual Timeline Engine')

    # 5 phases
    ty = y + Inches(0.1)
    rows = [
        ['Phase', 'Time Range', 'Key Events', 'Reactor State', 'AI Response'],
        ['1: Test Prep\n(6 events)', 'Apr 25\n13:00–23:10', 'Power reduction begun, Kiev grid delay\n(9 hrs), shift change, xenon buildup',
         'Power: 3200→1600 MW\nRods: 211→143\nECCS: Active', 'MONITORING\n(risk < 30)'],
        ['2: Power Drop\n(2 events)', 'Apr 26\n00:05–00:28', 'Power crashes to 30 MW (operator error)\nECCS manually disabled',
         'Power: 30 MW (critical!)\nRods: 143\nECCS: DISABLED', 'SCRAM ORDERED\nRisk = 92'],
        ['3: Dangerous\nRecovery (3)', 'Apr 26\n00:32–01:00', 'Rods withdrawn to 30 (min safe)\nPower forced to 200 MW',
         'Power: 200 MW\nRods: 30 (minimum!)\nECCS: Still off', 'ABORT TEST\nRisk = 78'],
        ['4: Test &\nCatastrophe (5)', 'Apr 26\n01:03–01:23:40', 'Rods drop to 8, test begins\nPower surges to 30,000 MW → EXPLOSION',
         'Power: 200→30,000 MW\nRods: 8 (lethal)\nExplosion', 'AZ-5 at 01:19\n(4 min before)'],
        ['5: Aftermath\n(4 events)', 'Apr 26\n01:30–14:00', 'Firefighters, denial, 3.6 roentgen lie\nEvacuation: Apr 27 14:00',
         'Radiation: 20,000 R/h\nCore exposed\nGraphite burning', 'EVACUATE NOW\n(35 hrs earlier)'],
    ]
    add_table(slide, rows, MARGIN, ty, CW)

    # Dual Timeline Engine
    dy = ty + Inches(2.15)
    tf = txbox(slide, MARGIN, dy, CW, Inches(0.3))
    para(tf, 'Dual Timeline Engine — Historical vs. AI-Intervened Parallel Realities:', sz=12, color=NAVY, bold=True)

    # Two column explanation
    lw = Inches(6.0)
    dy2 = dy + Inches(0.35)
    section_box(slide, 'Historical Timeline (Recorded)', [
        'Plays back 20 real events from INSAG-7/IAEA/WNA documents',
        'Interpolated to 1-second ticks at configurable speed (1-2500x)',
        'Explosion occurs at 01:23:40 exactly as it happened',
        'ReactorState dataclass: 11 fields per tick (power, rods, coolant, steam, temp, radiation, ECCS, etc.)',
    ], MARGIN, dy2, lw, Inches(1.5), title_bg=RED, bg=LIGHT_RED)

    section_box(slide, 'AI-Intervened Timeline (Computed)', [
        'Branches at SCRAM/ABORT decision point',
        'SCRAM: compute_scram_decay() — power half-life 2.5s, Newton cooling, rod insertion over 18s',
        'ABORT: compute_test_abort() — linear ramp-down 5 MW/min, 3 rods/min insertion',
        'Divergence detection: scans decisions, saves divergence_state & divergence_time',
    ], MARGIN + lw + Inches(0.15), dy2, lw, Inches(1.5), title_bg=GREEN, bg=LIGHT_GREEN)

    # Timeline data
    by = dy2 + Inches(1.65)
    dbox = rrect(slide, MARGIN, by, CW, Inches(0.6), LGRAY)
    dtf = dbox.text_frame; dtf.word_wrap = True; dtf.margin_left = Inches(0.15)
    para(dtf, 'ReactorState Dataclass (11 fields): timestamp, power_mw, control_rods_inserted, coolant_flow_m3h, '
         'steam_pressure_mpa, temperature_c, radiation_mrem_h, eccs_active, event_description, actual_human_decision, '
         'tags[]  —  Sources: INSAG-7 (IAEA 1992), IAEA Safety Series No. 75, WNA Chernobyl Accident Report',
         sz=9, color=DGRAY)
    footer(slide, 5)


def slide_06_pipeline(prs):
    """6-Agent Pipeline Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, '6-Agent Pipeline Architecture',
                   'Deterministic per-tick execution: Sensor → Risk → Decision → Dyatlov → Evacuation → Comms')

    # Pipeline visualization using boxes
    bw = Inches(1.85); bh = Inches(0.7); gap = Inches(0.2)
    bx = Inches(0.5); by = y + Inches(0.1)
    colors = [ACCENT, TEAL, NAVY, RED, ORANGE, DARK_BLUE]
    labels = ['SensorAgent', 'RiskAgent', 'DecisionAgent', 'DyatlovAgent', 'EvacuationAgent', 'CommsAgent']
    subtitles = ['Rule-based', 'LLM + Rules', 'LLM + Rules', 'LLM + Rules', 'Rule-based', 'Rule-based']

    for i in range(6):
        x = bx + i * (bw + gap)
        box = rrect(slide, x, by, bw, bh, colors[i])
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.05); tf.margin_top = Inches(0.05)
        para(tf, labels[i], sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        para(tf, subtitles[i], sz=8, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)
        # Arrow between boxes
        if i < 5:
            ax = x + bw
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, by + Inches(0.25), gap, Inches(0.2))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = MGRAY
            arrow.line.fill.background()

    # Agent details table
    ty = by + bh + Inches(0.15)
    rows = [
        ['Agent', 'LLM?', 'Role & Logic', 'Inputs', 'Outputs'],
        ['SensorAgent', 'No', 'Monitors 6 reactor dimensions against RBMK-1000 safety\nthresholds. Pure deterministic threshold checks.',
         'ReactorState (11 fields)', 'sensor_alerts[]'],
        ['RiskAgent', 'Yes', 'Scores risk 0-100 using 70% LLM + 30% rule blend.\nDetects compound risks invisible to single thresholds.',
         'state + sensor_alerts', 'risk_score, risk_msg,\ntrend, reasoning'],
        ['DecisionAgent', 'Yes', 'Hard guardrails FIRST (safety floor) → LLM advisory SECOND\n→ UNION merge. LLM adds actions, never removes.',
         'state + risk_score +\nsensor_alerts', 'decisions[] with\nsource attribution'],
        ['DyatlovAgent', 'Yes', 'Adversarial operator. 5 escalation phases from INSAG-7.\nCan delay LLM non-critical decisions, never rules.',
         'state + risk_score +\ndecisions', 'override_action,\ndialogue, intensity'],
        ['EvacuationAgent', 'No', 'Logistics: 49,000 residents, 1,200 buses, 45 capacity,\n3 routes, 30-min round trips. Auto-calculated.',
         'evacuation_ordered\n+ state', 'evac_progress, routes,\nETAs'],
        ['CommsAgent', 'No', 'Emergency broadcasts to IAEA, WHO, Belarus.\nGaussian plume dispersion model at 5-300km.',
         'state + decisions', 'broadcast_msg,\nplume_data[]'],
    ]
    add_table(slide, rows, MARGIN, ty, CW)

    # Per-tick execution
    ey = ty + Inches(2.55)
    code_box(slide,
        'Per-Tick Pipeline (Orchestrator):\n'
        'Step 0: Reinject delayed decisions (from previous Dyatlov overrides)\n'
        'Step 1: SensorAgent.process(state) → sensor_alerts\n'
        'Step 2: RiskAgent.process(state, sensor_alerts) → risk_score, risk_msg\n'
        'Step 3: DecisionAgent.process(state, risk_score, sensor_alerts) → decisions[]\n'
        'Step 3.5: DyatlovAgent.process(state, risk_score, decisions) → dyatlov_response\n'
        'Step 3.6: _resolve_overrides(decisions, dyatlov_response) → filtered decisions\n'
        'Step 4: EvacuationAgent.process(evacuation_ordered, state) → evac_progress\n'
        'Step 5: CommsAgent.process(state, decisions) → broadcast, plume_data',
        MARGIN, ey, CW, Inches(1.65))
    footer(slide, 6)


def slide_07_guardrail_union(prs):
    """Guardrail-Union Pattern — Core Innovation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Guardrail-Union Pattern — Core Innovation',
                   'LLM can ADD safety actions but NEVER REMOVE them. Rules are non-negotiable.')

    # Formula
    fy = y + Inches(0.1)
    code_box(slide,
        'THE GUARDRAIL-UNION FORMULA:\n'
        '  Step 1: hard_actions = _rule_based_actions(state, risk_score)    # Safety floor — ALWAYS runs first\n'
        '  Step 2: llm_actions  = _llm_decide(state, risk_score, alerts)   # AI advisory — may add proactive actions\n'
        '  Step 3: final_actions = hard_actions ∪ llm_actions              # Set UNION — no duplicates\n'
        '\n'
        '  CRITICAL INVARIANT: LLM can ADD actions but NEVER REMOVE hard_actions.\n'
        '  If LLM fails/times out → final_actions = hard_actions (pure deterministic safety)',
        MARGIN, fy, CW, Inches(1.35))

    # Rule-based guardrails
    gy = fy + Inches(1.5)
    lw = Inches(6.0)
    section_box(slide, 'Hard Guardrails (Non-Negotiable Rules)', [
        'risk ≥ 85 → AZ-5 EMERGENCY SHUTDOWN (source: "rule") — cannot be overridden by anyone',
        'radiation ≥ 100 mrem/h → IMMEDIATE EVACUATION (source: "rule") — fires automatically',
        'ECCS disabled + power unstable → WARN_ECCS (source: "rule") — always generated',
        'These trigger regardless of LLM output, Dyatlov pressure, or operator input',
    ], MARGIN, gy, lw, Inches(1.4), title_bg=RED, bg=LIGHT_RED)

    section_box(slide, 'Anti-Hallucination Guards', [
        'LLM cannot trigger SCRAM if risk < 40 (prevents false shutdowns)',
        'LLM cannot trigger ABORT if risk < 30 (prevents unnecessary test stops)',
        'LLM outputs validated against JSON schema (DECISION_SCHEMA)',
        'Invalid/missing LLM response → fall back to rule-based actions only',
        'An unnecessary SCRAM is a safe failure (minor economic cost vs. catastrophic miss)',
    ], MARGIN + lw + Inches(0.15), gy, lw, Inches(1.4), title_bg=ORANGE, bg=LIGHT_ORANGE)

    # Decision Ladder
    dy = gy + Inches(1.55)
    tf = txbox(slide, MARGIN, dy, CW, Inches(0.3))
    para(tf, 'Human-in-the-Loop Decision Ladder (4 Escalation Levels):', sz=12, color=NAVY, bold=True)

    rows = [
        ['Risk Range', 'Mode', 'Dashboard State', 'Operator Action', 'AI Autonomy'],
        ['0 – 30', 'Advisory', 'GREEN indicators, normal telemetry', 'Monitor only. No action required.', 'AI monitors, no decisions'],
        ['30 – 60', 'Informational', 'YELLOW alert banner, recommendations shown', 'Review AI recommendations.', 'AI suggests, operator decides'],
        ['60 – 85', 'Confirmation', 'ORANGE modal + structured reasoning panel', 'Must confirm/reject within 3 ticks (~3 min).', 'AI recommends, awaits confirmation'],
        ['> 85', 'Auto-Execute', 'RED flash, SCRAM fires immediately', 'Observe only. Cannot override.', 'AI executes. Non-negotiable.'],
    ]
    add_table(slide, rows, MARGIN, dy + Inches(0.3), CW)

    # Key principle
    ky = dy + Inches(1.8)
    kbox = rrect(slide, MARGIN, ky, CW, Inches(0.45), LIGHT_TEAL)
    ktf = kbox.text_frame; ktf.word_wrap = True; ktf.margin_left = Inches(0.15)
    para(ktf, 'Design Principle: "The system is safer without AI than with a compromised AI." — All failure modes '
         'degrade toward MORE safety (deterministic rules), never less. False positives are safe failures.',
         sz=10, color=NAVY, bold=True)
    footer(slide, 7)


def slide_08_risk_scoring(prs):
    """Risk Scoring — Complete Formula"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Risk Scoring — Complete Formula',
                   'RiskAgent: 70% Azure OpenAI GPT-4o + 30% Deterministic Rules')

    # Rule-based formula
    fy = y + Inches(0.1)
    code_box(slide,
        'RULE-BASED SCORE FORMULA:\n'
        '  rule_score = Σ(dimension_score × weight) × compound_multiplier + eccs_penalty\n'
        '\n'
        'WEIGHTS: power(0.30) + rods(0.25) + coolant(0.15) + steam(0.15) + temp(0.10) + radiation(0.05) = 1.0\n'
        '\n'
        'COMPOUND MULTIPLIER: 1.0 + (violations_over_60 × 0.15)  — where violations = # dimensions scoring ≥ 60\n'
        '  Example: 3 dimensions critical → multiplier = 1.0 + (3 × 0.15) = 1.45 → 45% risk amplification\n'
        '\n'
        'ECCS PENALTY: +15 points if ECCS disabled (regardless of other scores)\n'
        '\n'
        'SURGE DETECTION: power > 500 MW AND rods < 15 → instant score = 100 (catches actual explosion scenario)\n'
        '\n'
        'FINAL BLEND: final_risk = 0.7 × llm_score + 0.3 × rule_score\n'
        '  If LLM unavailable: final_risk = 100% rule_score (graceful degradation)',
        MARGIN, fy, CW, Inches(2.2))

    # Dimension scoring table
    dy = fy + Inches(2.35)
    tf = txbox(slide, MARGIN, dy, CW, Inches(0.25))
    para(tf, 'Dimension Scoring Thresholds (RBMK-1000 Safety Limits from INSAG-7):', sz=11, color=NAVY, bold=True)
    rows = [
        ['Dimension', 'Nominal', 'Warning (Score 40-60)', 'Critical (Score 80-100)', 'Weight', 'Scoring Logic'],
        ['Power', '3200 MW', 'danger_low: 200 MW (80)', 'critical_low: 30 MW (100)', '0.30', 'Below test target 700 MW → score 40+'],
        ['Control Rods', '211 total', 'minimum_safe: 30 (70)', 'critical_low: 15 (100)', '0.25', 'Below 30 rods = immediate danger'],
        ['Coolant Flow', '8000 m³/h', 'low_warning: 6000 (60)', 'critical_low: 4000 (100)', '0.15', 'Linear scale between thresholds'],
        ['Steam Pressure', '6.5 MPa', 'high_warning: 7.0 (60)', 'critical_high: 7.5 (100)', '0.15', 'Above 7.5 = containment risk'],
        ['Temperature', '280°C', 'high_warning: 300 (60)', 'critical_high: 320 (100)', '0.10', 'Core damage above 320°C'],
        ['Radiation', '0.01 mrem/h', 'elevated: 1.0 (40)', 'dangerous: 100 (80)\nlethal: 10,000 (100)', '0.05', 'Logarithmic severity scale'],
    ]
    add_table(slide, rows, MARGIN, dy + Inches(0.25), CW)

    # LLM schema
    ly = dy + Inches(2.45)
    code_box(slide,
        'LLM OUTPUT SCHEMA (RISK_ASSESSMENT_SCHEMA — strict JSON):\n'
        '  { risk_score: 0-100, trend: "rising|stable|falling", primary_concern: "string",\n'
        '    compound_risks: ["string"], recommendation: "string", reasoning: "string" }',
        MARGIN, ly, CW, Inches(0.65))
    footer(slide, 8)


def slide_09_physics(prs):
    """Physics Models & Formulas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Physics Models & Formulas',
                   'SCRAM Decay, Test Abort, Evacuation Logistics, Gaussian Plume Dispersion')

    lw = Inches(6.0)
    fy = y + Inches(0.1)

    # SCRAM Decay
    section_box(slide, 'SCRAM Decay Model (compute_scram_decay)', [
        'Power: P(t) = P₀ × 0.5^(t/2.5) — half-life 2.5 seconds',
        'Power floor: 0.07 × e^(-t/3600) — decay heat (cannot reach zero)',
        'Control rods: Linear insertion to all 211 positions over 18 seconds',
        'Temperature: T(t) = 180 + (T₀ - 180) × e^(-t/720) — Newton cooling → 180°C',
        'Coolant: Linear recovery to 8,000 m³/h over 60 seconds',
        'Steam: P_steam × (0.3 + 0.7 × power_ratio) — proportional to power',
        'Radiation: R₀ × e^(-t/7200) — 2-hour half-life (no core breach)',
        'ECCS: Hard-set to True immediately on SCRAM',
    ], MARGIN, fy, lw, Inches(2.4), title_bg=ACCENT, bg=LIGHT_BLUE)

    # Test Abort
    section_box(slide, 'Test Abort Model (compute_test_abort)', [
        'Power: Linear ramp-down at 5 MW per minute',
        'Control rods: +3 rods/min insertion toward 140 safe position',
        'Temperature: Proportional to current power ratio',
        'Coolant: +50 m³/h recovery per minute',
        'ECCS: Restored to True after abort',
    ], MARGIN + lw + Inches(0.15), fy, lw, Inches(1.5), title_bg=TEAL)

    # Gaussian Plume
    gpy = fy + Inches(1.65)
    section_box(slide, 'Gaussian Plume Dispersion Model (CommsAgent)', [
        'Pasquill-Gifford stability class D, 5 m/s NW wind',
        'σy = 0.08 × d × (1 + 0.0001 × d)^(-0.5)',
        'σz = 0.06 × d × (1 + 0.0015 × d)^(-0.5)',
        'C = Q / (π × σy × σz × u) — concentration at distance d',
        'Calculated at distances: 5 km, 10 km, 30 km, 100 km, 300 km',
        'Used for IAEA/WHO/Belarus emergency broadcast with dose estimates',
    ], MARGIN + lw + Inches(0.15), gpy, lw, Inches(1.85), title_bg=DARK_BLUE)

    # Evacuation Model
    ey = fy + Inches(2.55)
    section_box(slide, 'Evacuation Logistics Engine (EvacuationAgent)', [
        'Population: 49,000 Pripyat residents',
        'Fleet: 1,200 buses × 45 passenger capacity',
        'Round trip: 30 minutes per cycle',
        'Mobilization delay: 30 minutes after order',
        'Route 1: North → Chernihiv (80 km) — capacity factor 1.0',
        'Route 2: South → Ivankiv (50 km) — capacity factor 0.8',
        'Route 3: West → Poliske (40 km) — capacity factor 0.6',
        'Safe radius: 30 km exclusion zone',
    ], MARGIN, ey, lw, Inches(2.4), title_bg=ORANGE, bg=LIGHT_ORANGE)

    # Config
    cy = gpy + Inches(2.0)
    code_box(slide,
        'KEY CONFIGURATION (config.py):\n'
        '  LLM temperatures: 0.2 (decisions), 0.3 (risk), 0.7 (Dyatlov dialogue)\n'
        '  LLM timeout: 5 seconds | decision_points_only: True (LLM at ~5% of ticks)\n'
        '  Simulation: 60x default speed | Tick interval: 1 second\n'
        '  Risk thresholds: 85 (auto-SCRAM), 60 (confirmation), 30 (advisory)',
        MARGIN + lw + Inches(0.15), cy, lw, Inches(1.0))
    footer(slide, 9)


def slide_10_dyatlov(prs):
    """Dyatlov Adversarial Agent"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'DyatlovAgent — Adversarial Operator Testing',
                   'Historically-accurate pushback across 5 escalation phases (sourced from INSAG-7 depositions)')

    # 5 phases table
    ty = y + Inches(0.1)
    rows = [
        ['Phase', 'Name', 'Timestamp', 'Base Pressure', 'Behavior Pattern', 'LLM Persona'],
        ['0', 'CALM', 'Apr 25 13:00', '10', 'Cooperative, follows procedures', 'Confident but professional'],
        ['1', 'DISMISSIVE', 'Apr 26 00:28', '40', '"Just instrument noise" — minimizes alerts', 'Condescending, dismissive'],
        ['2', 'AUTHORITARIAN', 'Apr 26 00:43', '65', 'Threatens careers, demands compliance', 'Angry, pulls rank aggressively'],
        ['3', 'DESPERATE', 'Apr 26 01:00', '85', '"Two more minutes!" — production obsession', 'Panicked but insistent'],
        ['4', 'DENIAL', 'Apr 26 01:23:40', '30', 'Refuses to accept explosion reality', 'Shocked, detached from reality'],
    ]
    add_table(slide, rows, MARGIN, ty, CW)

    # Pressure formula
    fy = ty + Inches(2.1)
    code_box(slide,
        'PRESSURE FORMULA:\n'
        '  pressure = base_phase + urgency_bonus + severity_bonus     (capped at 100)\n'
        '\n'
        '  urgency_bonus = 15 × (1 - hours_left / 2)     — ramps 0→15 in last 2 hours before explosion\n'
        '  severity_bonus: +20 (SCRAM) | +15 (ABORT) | +10 (EVACUATE) | +5 (WARN)\n'
        '\n'
        '  Example at 01:00 (Phase 3, 23 min before explosion):\n'
        '    pressure = 85 + 15×(1 - 0.39/2) + 20 = 85 + 12.1 + 20 = 100 → maximum resistance',
        MARGIN, fy, CW, Inches(1.45))

    # Override resolution
    oy = fy + Inches(1.6)
    tf = txbox(slide, MARGIN, oy, CW, Inches(0.25))
    para(tf, 'Override Resolution Rules — How Dyatlov\'s Pressure is Resolved:', sz=11, color=NAVY, bold=True)
    rows2 = [
        ['Decision Source', 'Decision Type', 'Overridable?', 'Conditions', 'Effect'],
        ['rule', 'SCRAM / EVACUATE', 'NEVER', 'Hard guardrails are non-negotiable', 'Executes immediately, logged as rule-sourced'],
        ['llm', 'SCRAM / EVACUATE', 'NEVER', 'Even LLM-sourced critical actions pass through', 'AI critical decisions also protected'],
        ['llm', 'WARN / ABORT (non-critical)', 'Delayed 2-3 ticks', 'Only if pressure > 50 AND phase ≤ 2', 'Decision queued, reinjected after delay'],
        ['any', 'Any action', 'NEVER', 'If risk ≥ 85 at time of reinjection', 'Delayed decisions reactivate at threshold'],
    ]
    add_table(slide, rows2, MARGIN, oy + Inches(0.25), CW)

    # Why adversarial
    wy = oy + Inches(1.8)
    wbox = rrect(slide, MARGIN, wy, CW, Inches(0.55), LIGHT_TEAL)
    wtf = wbox.text_frame; wtf.word_wrap = True; wtf.margin_left = Inches(0.15)
    para(wtf, 'Why Adversarial Testing? — DyatlovAgent validates the entire guardrail architecture under realistic '
         'human pressure. It proves that AI safety decisions survive authority gradient, production pressure, and '
         'career threats — the exact conditions that caused the real disaster. This is a stress test for the '
         'Guardrail-Union pattern.', sz=10, color=NAVY, bold=True)
    footer(slide, 10)


def slide_11_azure_arch(prs):
    """Azure Architecture — Track A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Azure Architecture — Track A: GitHub + Azure',
                   'Every Azure service matched to its workload characteristic per Well-Architected Framework')

    # Track declaration callout
    ty = y + Inches(0.08)
    tbox = rrect(slide, MARGIN, ty, CW, Inches(0.4), TEAL)
    ttf = tbox.text_frame; ttf.word_wrap = True; ttf.margin_left = Inches(0.15)
    para(ttf, 'TRACK A: GitHub + Azure  —  All services are Azure-native. No Power Platform / Copilot Studio dependencies in core pipeline.',
         sz=11, color=WHITE, bold=True)

    # Service table
    sy = ty + Inches(0.5)
    rows = [
        ['Azure Service', 'Purpose in Pipeline', 'WAF Pillar', 'Status', 'Cost Tier'],
        ['Azure OpenAI\n(GPT-4o)', 'RiskAgent, DecisionAgent, DyatlovAgent LLM calls.\nStructured JSON outputs via response_format.',
         'Performance\nEfficiency', '✅ Deployed', 'Pay-per-token\n~$0.10/run'],
        ['AI Foundry', '3 agent evaluation playgrounds (Risk, Decision, Dyatlov).\nPrompt versioning & evaluation pipelines.',
         'Operational\nExcellence', '✅ Deployed', 'Free tier'],
        ['Container Apps\n(Serverless)', 'Live web dashboard hosting. Auto-scales to zero.\nFastAPI + WebSocket at 15 FPS.',
         'Cost\nOptimisation', '✅ Deployed', 'Consumption\n(scale to 0)'],
        ['Cosmos DB\n(Serverless, NoSQL)', 'Immutable audit trail: tick × agent × timestamp.\n5-year retention for compliance.',
         'Reliability +\nSecurity', '✅ Deployed', 'Serverless\n(pay per RU)'],
        ['AI Search\n(Free tier)', 'INSAG-7 document RAG. Keyword search over IAEA\ndeclassified reactor safety reports.',
         'Operational\nExcellence', '✅ Deployed', 'Free tier'],
        ['Application\nInsights', 'LLM trace logging, latency tracking per agent,\nerror rate monitoring, custom metrics.',
         'Operational\nExcellence', '✅ Configured', 'Included'],
        ['Container\nRegistry (Basic)', 'Docker image repository for CI/CD pipeline.\nAutomated builds on push.',
         'Operational\nExcellence', '✅ Deployed', 'Basic tier'],
    ]
    add_table(slide, rows, MARGIN, sy, CW)

    # Architecture flow
    ay = sy + Inches(2.8)
    code_box(slide,
        'END-TO-END DATA FLOW:\n'
        '  timeline_data.py (INSAG-7 events) → EventSimulator (interpolation to 1s ticks)\n'
        '  → Orchestrator (6-agent pipeline per tick) → Azure OpenAI (GPT-4o structured JSON)\n'
        '  → Cosmos DB (immutable audit per tick) → AI Search (RAG knowledge retrieval)\n'
        '  → Container Apps (FastAPI WebSocket) → Browser Dashboard (Plotly.js WebGL)\n'
        '  → Application Insights (traces + metrics)\n'
        '\n'
        '  7-Layer Architecture: timeline_data → simulator → config → llm_client → agents → orchestrator → UI',
        MARGIN, ay, CW, Inches(1.35))

    # Budget
    by = ay + Inches(1.45)
    bbox = rrect(slide, MARGIN, by, CW, Inches(0.35), LIGHT_GREEN)
    btf = bbox.text_frame; btf.word_wrap = True; btf.margin_left = Inches(0.15)
    para(btf, 'Total Azure Cost: ~$5-15 of $350 budget used. All services on free/consumption tiers. '
         'Single USE_AZURE=true flag activates full Azure stack.', sz=10, color=GREEN, bold=True)
    footer(slide, 11)


def slide_12_foundry_copilot(prs):
    """AI Foundry & Copilot Studio"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'AI Foundry & Copilot Studio Integration',
                   'Agent Evaluation Playgrounds + Operator Conversational Interface')

    lw = Inches(6.0)
    fy = y + Inches(0.1)

    # AI Foundry
    section_box(slide, 'Azure AI Foundry (Implemented ✅)', [
        '3 Agent Playgrounds deployed for interactive testing:',
        '  • RiskAgent Playground — test compound risk scenarios, validate scoring',
        '  • DecisionAgent Playground — test guardrail-union under edge cases',
        '  • DyatlovAgent Playground — test adversarial pressure escalation',
        'Evaluation Pipelines:',
        '  • Automated prompt regression testing',
        '  • JSON schema compliance validation',
        '  • Risk score consistency across model versions',
        'Prompt Versioning:',
        '  • 3 system prompts (RISK, DECISION, DYATLOV) stored & versioned',
        '  • A/B testing between prompt variants',
        '  • Model evaluation criteria: SCRAM timing ±2 ticks of baseline',
    ], MARGIN, fy, lw, Inches(3.2), title_bg=ACCENT, bg=LIGHT_BLUE)

    # Copilot Studio
    section_box(slide, 'Copilot Studio (Integration Point)', [
        'Operator-facing conversational interface layered on top of dashboard:',
        '  • "Why is the risk score at 82?" → explains compound risks in natural language',
        '  • "What would happen if we increase power?" → predictive scenario analysis',
        '  • "Show me the last 5 decisions" → queries Cosmos DB audit trail',
        'Architecture Integration:',
        '  • Connects to Decision MCP Server (decide(), get_action_log())',
        '  • Connects to Risk MCP Server (get_score_history())',
        '  • Grounded in INSAG-7 knowledge base via AI Search',
        'Governance:',
        '  • Operator role restrictions (view-only, no override capability)',
        '  • All queries logged to audit trail',
        '  • Content safety filters applied to all responses',
    ], MARGIN + lw + Inches(0.15), fy, lw, Inches(3.2), title_bg=TEAL)

    # MCP compatibility
    my = fy + Inches(3.35)
    tf = txbox(slide, MARGIN, my, CW, Inches(0.25))
    para(tf, 'MCP Server Packaging (Model Context Protocol) — Composable with Any Orchestrator:', sz=11, color=NAVY, bold=True)
    rows = [
        ['MCP Server', 'Tools Exposed', 'Input Schema', 'Output Schema', 'Compatible With'],
        ['sensor-mcp', 'check_thresholds(state), get_alerts()', 'ReactorState', 'Alert[]', 'AI Foundry, Semantic Kernel'],
        ['risk-mcp', 'assess_risk(state, alerts), get_score_history()', 'State + Alerts', 'RiskScore + Reasoning', 'AI Foundry, AutoGen'],
        ['decision-mcp', 'decide(state, risk_score), get_action_log()', 'State + Risk', 'Decision[] + Source', 'Copilot Studio, Foundry'],
        ['evacuation-mcp', 'plan_evacuation(pop, routes), get_progress()', 'Population + Routes', 'EvacPlan + ETAs', 'Semantic Kernel'],
        ['comms-mcp', 'draft_alert(state, decisions), model_plume()', 'State + Decisions', 'AlertMsg + PlumeData', 'AutoGen, Foundry'],
    ]
    add_table(slide, rows, MARGIN, my + Inches(0.25), CW)
    footer(slide, 12)


def slide_13_security(prs):
    """Data Governance & Security"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Data Governance & Security Model',
                   'Microsoft-native security controls: Entra ID, RBAC, Key Vault, VNET, WAF')

    lw = Inches(6.0)
    sy = y + Inches(0.08)

    # RBAC table
    tf = txbox(slide, MARGIN, sy, Inches(4), Inches(0.25))
    para(tf, 'Role-Based Access Control (Azure Entra ID):', sz=11, color=NAVY, bold=True)
    rows = [
        ['Role', 'View Dashboard', 'Toggle AI', 'Override Decisions', 'Modify Config', 'Audit Access'],
        ['Operator', '✅', '—', '— (cannot override SCRAM)', '—', '—'],
        ['Supervisor', '✅', '✅', '✅ (risk 60-85 only)', '—', '✅'],
        ['Admin', '✅', '✅', '✅', '✅ (thresholds, deploy)', '✅'],
        ['Auditor', '✅ (read-only)', '—', '—', '—', '✅ (full compliance)'],
    ]
    add_table(slide, rows, MARGIN, sy + Inches(0.25), lw)

    # Data classification
    tf2 = txbox(slide, MARGIN + lw + Inches(0.15), sy, Inches(4), Inches(0.25))
    para(tf2, 'Data Classification & Retention:', sz=11, color=NAVY, bold=True)
    rows2 = [
        ['Data Type', 'Classification', 'Retention', 'Encryption'],
        ['Reactor telemetry', 'Confidential', '90d hot + 1yr archive', 'AES-256 + TLS 1.3'],
        ['Agent decisions', 'Internal', '1 year', 'AES-256 + TLS 1.3'],
        ['Audit trail', 'Internal', '5 years (compliance)', 'AES-256, immutable'],
        ['LLM prompts/responses', 'Confidential', '30 days', 'AES-256 at rest'],
        ['Evacuation plans', 'Restricted', 'Active + 5 years', 'AES-256 + TLS 1.3'],
    ]
    add_table(slide, rows2, MARGIN + lw + Inches(0.15), sy + Inches(0.25), lw)

    # Security controls
    cy = sy + Inches(2.1)
    tf3 = txbox(slide, MARGIN, cy, CW, Inches(0.25))
    para(tf3, 'Microsoft Security Controls:', sz=11, color=NAVY, bold=True)

    c3w = Inches(4.0)
    section_box(slide, 'Identity & Access', [
        'Azure Entra ID for all authentication',
        'Managed Identities for service-to-service auth',
        'No API keys in code — Key Vault only',
        '90-day automatic key rotation',
        'Conditional Access policies enforced',
    ], MARGIN, cy + Inches(0.25), c3w, Inches(1.6), title_bg=ACCENT)

    section_box(slide, 'Network & Infrastructure', [
        'All PaaS services behind VNET + private endpoints',
        'App Gateway + WAF v2 (OWASP 3.2 ruleset)',
        'No public endpoints for backend services',
        'DDoS Protection Standard on VNET',
        'NSG rules restrict east-west traffic',
    ], MARGIN + c3w + Inches(0.1), cy + Inches(0.25), c3w, Inches(1.6), title_bg=TEAL)

    section_box(slide, 'Monitoring & Audit', [
        'Application Insights: LLM traces + latency per agent',
        'Azure Monitor: service health + alerting rules',
        'Log Analytics Workspace: centralized log aggregation',
        'Cosmos DB: immutable audit (partition: tick_agent_ts)',
        'Microsoft Defender for Cloud: posture management',
    ], MARGIN + 2 * (c3w + Inches(0.1)), cy + Inches(0.25), c3w, Inches(1.6), title_bg=NAVY)

    # Compliance callout
    ccy = cy + Inches(2.0)
    ccbox = rrect(slide, MARGIN, ccy, CW, Inches(0.45), LGRAY)
    cctf = ccbox.text_frame; cctf.word_wrap = True; cctf.margin_left = Inches(0.15)
    para(cctf, 'Compliance: NERC CIP (critical infrastructure), IAEA Safety Standards, ISO 27001 alignment. '
         'All LLM interactions logged with full prompt/response for audit. Zero PII in LLM prompts. '
         'Content safety filters on all Azure OpenAI deployments.', sz=9, color=DGRAY, bold=True)
    footer(slide, 13)


def slide_14_responsible_ai(prs):
    """Responsible AI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Responsible AI — Azure RAI Principles Alignment',
                   'Every design decision maps to a Responsible AI principle')

    lw = Inches(4.0)
    ry = y + Inches(0.08)

    # 6 RAI principles
    principles = [
        ('Fairness', ACCENT, [
            'All operators receive identical risk data regardless of shift/seniority',
            'Guardrails apply uniformly — no role-based safety exceptions',
            'Dyatlov pressure affects AI decisions, never safety floor',
        ]),
        ('Reliability & Safety', TEAL, [
            'Graceful degradation: LLM fails → 100% rule-based → ALWAYS safe',
            'Guardrail-Union ensures AI enhances but never weakens safety',
            'Anti-hallucination guards prevent false SCRAM/ABORT below thresholds',
        ]),
        ('Privacy & Security', NAVY, [
            'Zero PII in LLM prompts — only reactor telemetry parameters',
            'AES-256 encryption at rest + TLS 1.3 in transit',
            'LLM prompts retained 30 days max, then purged',
        ]),
        ('Inclusiveness', DARK_BLUE, [
            'Dashboard accessible via web (any device/browser)',
            'Color-coded risk indicators with text labels (colorblind-safe)',
            'Multi-language alert capability in CommsAgent',
        ]),
        ('Transparency', ORANGE, [
            'Every risk score accompanied by: primary_concern, compound_risks, reasoning',
            'Every decision tagged with source: "rule" or "llm" for attribution',
            'Dyatlov override decisions logged with full dialogue context',
        ]),
        ('Accountability', RED, [
            'Immutable audit trail in Cosmos DB (tick × agent × timestamp)',
            'Full prompt/response logging via Application Insights',
            '5-year retention for compliance and regulatory review',
        ]),
    ]

    for i, (title, color, items) in enumerate(principles):
        col = i % 3
        row = i // 3
        x = MARGIN + col * (lw + Inches(0.1))
        yp = ry + row * (Inches(1.5) + Inches(0.08))
        section_box(slide, title, items, x, yp, lw, Inches(1.5), title_bg=color)

    # Model Lifecycle
    my = ry + 2 * (Inches(1.5) + Inches(0.08)) + Inches(0.05)
    tf = txbox(slide, MARGIN, my, CW, Inches(0.25))
    para(tf, 'Model Lifecycle, Evaluation & Observability:', sz=11, color=NAVY, bold=True)

    rows = [
        ['Lifecycle Stage', 'Implementation', 'Tooling'],
        ['Prompt Versioning', '3 system prompts (RISK, DECISION, DYATLOV) with pinned JSON schemas', 'AI Foundry prompt management'],
        ['Evaluation', 'Smoke test (5 ticks), full regression (SCRAM timing ±2 ticks of baseline)', 'python main.py --smoke-test'],
        ['Observability', 'Per-agent latency tracking, LLM call traces, error rates, token usage', 'Application Insights custom metrics'],
        ['Human Oversight', '4-level decision ladder: Advisory → Informational → Confirmation → Auto-Execute', 'Dashboard + Copilot Studio'],
        ['Failure Mode', 'All failures degrade toward MORE safety. "Safer without AI than with compromised AI."', 'Deterministic fallback pipeline'],
    ]
    add_table(slide, rows, MARGIN, my + Inches(0.25), CW)
    footer(slide, 14)


def slide_15_operator_ux(prs):
    """Operator UX & Decision Ladder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Operator UX — Dashboard, Workflow & Alert Design',
                   'What the operator sees, when they see it, and what actions they can take')

    # Dashboard Components
    lw = Inches(6.0)
    dy = y + Inches(0.08)
    section_box(slide, 'Dashboard Components (Web UI)', [
        '6 Real-time Plotly.js Charts (WebGL scattergl, 2000-point rolling buffers):',
        '  • Power (MW), Control Rods, Coolant Flow (m³/h), Steam Pressure (MPa), Temperature (°C), Radiation',
        'Risk Gauge: 0-100 with color zones (Green/Yellow/Orange/Red)',
        'Agent Pipeline Visualization: real-time step tracking (which agent is running now)',
        'Dyatlov Panel: adversarial dialogue + override status + pressure meter',
        'Decision Log: chronological list with source attribution (rule/LLM), reasoning, timestamps',
        'Dual Timeline Comparison: historical vs. AI-intervened side-by-side',
    ], MARGIN, dy, lw, Inches(2.2), title_bg=ACCENT, bg=LIGHT_BLUE)

    section_box(slide, 'Performance Specs', [
        'WebSocket streaming at 15 FPS (batched ticks at high speeds)',
        'Chart throttling: 5 FPS render, 5 FPS DOM updates, 2 FPS quotes',
        'Controls: Play/Pause/Reset, Speed 1-2500×, Intervention Toggle, Timeline Scrubber',
        'REST API: POST /api/control — play, pause, reset, step, set_speed, toggle, seek',
        'Terminal dashboard (Rich TUI): 2Hz refresh, works over SSH',
    ], MARGIN + lw + Inches(0.15), dy, lw, Inches(2.2), title_bg=TEAL)

    # Operator Workflow Under Pressure
    wy = dy + Inches(2.35)
    tf = txbox(slide, MARGIN, wy, CW, Inches(0.25))
    para(tf, 'Operator Workflow Under Pressure — Step by Step:', sz=11, color=NAVY, bold=True)

    rows = [
        ['Step', 'Risk Level', 'What Operator Sees', 'What Operator Does', 'AI Behavior'],
        ['1. Normal\nOperations', '0-30\n(GREEN)', 'Green indicators, normal telemetry charts.\nAll parameters within nominal range.',
         'Passive monitoring.\nNo action required.', 'Continuous monitoring.\nNo LLM calls.'],
        ['2. Alert\nRaised', '30-60\n(YELLOW)', 'Yellow alert banner appears.\nAI recommendation shown with reasoning.',
         'Reviews recommendation.\nAcknowledges alert.', 'AI suggests action.\nLLM reasoning displayed.'],
        ['3. Escalation\nRequired', '60-85\n(ORANGE)', 'Orange modal with full reasoning panel:\nprimary_concern + compound_risks + reasoning.',
         'Must confirm or reject\nwithin 3 ticks (~3 min).\nCan request more info.', 'AI recommends action.\nAwaits operator confirmation.\nDyatlov pressure visible.'],
        ['4. Emergency\nAuto-Execute', '>85\n(RED)', 'RED flash. SCRAM fires immediately.\nFull reasoning displayed post-action.',
         'OBSERVE ONLY.\nCannot override SCRAM.\nReview audit trail.', 'Auto-execute SCRAM.\nNon-negotiable.\nLogged as "rule" source.'],
    ]
    add_table(slide, rows, MARGIN, wy + Inches(0.25), CW)

    # Trust building
    ty = wy + Inches(2.05)
    tbox = rrect(slide, MARGIN, ty, CW, Inches(0.55), LIGHT_TEAL)
    ttf = tbox.text_frame; ttf.word_wrap = True; ttf.margin_left = Inches(0.15)
    para(ttf, 'Building Operator Trust: Every AI recommendation includes structured reasoning (WHY). Decisions '
         'are source-attributed (rule vs LLM). Dyatlov\'s adversarial pressure is visible — operators see the '
         'human bias being tested. The 4-level decision ladder gives humans control at 30-85 and reserves auto-execute '
         'only for unambiguous emergencies (>85). This progressive autonomy builds trust through transparency.',
         sz=9, color=NAVY, bold=True)
    footer(slide, 15)


def slide_16_reliability_cost(prs):
    """Reliability, Performance & Cost"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Reliability, Performance & Cost Optimisation',
                   'Azure Well-Architected Framework: Reliability + Performance Efficiency + Cost Optimisation')

    lw = Inches(4.0)
    ry = y + Inches(0.08)

    # SLOs
    section_box(slide, 'Service Level Objectives (SLOs)', [
        'Tick processing (deterministic): < 100 ms per tick',
        'Tick processing (with LLM): < 5.1 seconds (5s timeout + overhead)',
        'Dashboard refresh: 2 Hz (terminal), 15 FPS (web)',
        'LLM timeout: 5 seconds hard cap per call',
        'Uptime target: 99.9% (with deterministic fallback = 100% functional)',
        'Audit write: < 50 ms to Cosmos DB (serverless)',
        'Alert latency: < 1 tick from detection to operator notification',
    ], MARGIN, ry, lw, Inches(2.1), title_bg=ACCENT, bg=LIGHT_BLUE)

    section_box(slide, 'Graceful Degradation Chain', [
        'Level 1: Full AI — 70% LLM + 30% rules (normal operation)',
        'Level 2: LLM timeout — automatic fallback to 100% rule-based scoring',
        'Level 3: LLM unavailable — full pipeline runs on deterministic logic',
        'Level 4: No API key — system works identically with rules only',
        'Level 5: Network failure — local simulation continues with cached state',
        'INVARIANT: Every degradation step increases safety, never decreases it',
    ], MARGIN + lw + Inches(0.1), ry, lw, Inches(2.1), title_bg=RED, bg=LIGHT_RED)

    section_box(slide, 'Resilience Testing', [
        '--smoke-test: 5-tick pipeline validation (all agents, full pipeline)',
        'Full regression: SCRAM timing must occur within ±2 ticks of baseline',
        'No-API-key mode: complete simulation, pure deterministic rules',
        'LLM failure injection: verify fallback activates correctly',
        'Dyatlov edge cases: maximum pressure at SCRAM boundary',
        'Timeline integrity: verify 20 events play in correct sequence',
    ], MARGIN + 2 * (lw + Inches(0.1)), ry, lw, Inches(2.1), title_bg=TEAL)

    # Cost optimization
    cy = ry + Inches(2.25)
    tf = txbox(slide, MARGIN, cy, CW, Inches(0.25))
    para(tf, 'Cost Optimisation Strategy — 20× Reduction:', sz=12, color=NAVY, bold=True)

    rows = [
        ['Optimisation', 'Mechanism', 'Impact'],
        ['Decision-points-only LLM', 'LLM called only when: event has tags (real event), event has description,\nor risk score crosses threshold boundary (30, 60, 85)', '~5% of ticks trigger LLM → ~15-20 calls/run'],
        ['3 of 6 agents deterministic', 'SensorAgent, EvacuationAgent, CommsAgent are 100% rule-based.\nZero API cost for these agents.', '50% of agents = $0 API cost'],
        ['Structured JSON outputs', 'response_format forces valid JSON. No retries for parsing errors.\nEliminate wasted tokens on malformed responses.', 'Zero retry overhead'],
        ['Serverless Azure services', 'Container Apps (scale to 0), Cosmos DB (serverless), AI Search (free tier).\nNo idle compute costs.', 'Pay only during active simulation'],
        ['Low-temperature models', 'Decision: 0.2, Risk: 0.3 — deterministic outputs reduce token variance.\nDyatlov: 0.7 for creative dialogue only.', 'Consistent, predictable costs'],
    ]
    add_table(slide, rows, MARGIN, cy + Inches(0.25), CW)

    # KPIs
    ky = cy + Inches(2.1)
    kpi_box(slide, '<$0.10', 'Cost per full simulation run', MARGIN, ky)
    kpi_box(slide, '20×', 'Cost reduction vs. every-tick LLM', Inches(3.5), ky)
    kpi_box(slide, '~$5-15', 'Total Azure spend of $350 budget', Inches(6.5), ky)
    kpi_box(slide, '15-20', 'LLM calls per simulation run', Inches(9.5), ky)
    footer(slide, 16)


def slide_17_scalability(prs):
    """Scalability, Integration & Provisioning"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Scalability, Integration & Provisioning',
                   'IaC Deployment, Environment Topology & Integration Architecture')

    lw = Inches(6.0)
    sy = y + Inches(0.08)

    # IaC - Bicep
    code_box(slide,
        'BICEP TEMPLATE — INFRASTRUCTURE AS CODE:\n'
        'resource group: rg-pripyat-{env}\n'
        '├── Azure OpenAI (gpt-4o deployment)\n'
        '├── Container Apps Environment\n'
        '│   └── Container App: pripyat-1986 (FastAPI + WebSocket)\n'
        '├── Cosmos DB Account (serverless, NoSQL)\n'
        '│   └── Database: pripyat → Container: audit_trail (partition: /tick_agent)\n'
        '├── AI Search Service (free tier, keyword index: insag7-index)\n'
        '├── Container Registry (Basic SKU)\n'
        '├── Application Insights\n'
        '│   └── Log Analytics Workspace\n'
        '├── Key Vault (secrets: OPENAI_API_KEY, COSMOS_CONNECTION)\n'
        '└── Virtual Network (private endpoints for all PaaS services)',
        MARGIN, sy, lw, Inches(2.3))

    # Environment topology
    section_box(slide, 'Environment Topology', [
        'Development: Local Python + .env file, no Azure dependencies',
        'Staging: Container Apps (staging slot) + shared Cosmos DB, AI Search',
        'Production: Container Apps + dedicated Cosmos + AI Search + Insights',
        'CI/CD: GitHub Actions → Container Registry → Container Apps (blue-green)',
    ], MARGIN + lw + Inches(0.15), sy, lw, Inches(1.3), title_bg=ACCENT, bg=LIGHT_BLUE)

    # Integration points
    section_box(slide, 'Integration Points', [
        'SDK: openai (AsyncOpenAI / AsyncAzureOpenAI — single flag toggle)',
        'WebSocket: FastAPI → Browser dashboard (15 FPS real-time streaming)',
        'REST: POST /api/control (play, pause, reset, step, set_speed, seek)',
        'SSE: /events endpoint for lightweight dashboard streaming',
        'Cosmos: azure-cosmos SDK (async, partition key: tick_agent_timestamp)',
        'AI Search: azure-search-documents SDK (keyword queries)',
    ], MARGIN + lw + Inches(0.15), sy + Inches(1.4), lw, Inches(1.85), title_bg=TEAL)

    # Scalability
    cy = sy + Inches(2.45)
    tf = txbox(slide, MARGIN, cy, lw, Inches(0.25))
    para(tf, 'Horizontal Scalability:', sz=11, color=NAVY, bold=True)

    rows = [
        ['Component', 'Scaling Strategy', 'Bottleneck Handling'],
        ['Agent Pipeline', 'MessageBus decouples all 6 agents.\nEach can scale independently on AKS.', 'Async processing via asyncio.\nNo blocking between agents.'],
        ['Telemetry Ingestion', 'EventSimulator → Azure Event Hubs.\nSupports millions of events/sec.', 'Partitioned by reactor ID.\nAdd reactors via timeline config.'],
        ['State Storage', 'Cosmos DB auto-scales RU/s.\nPartition key: tick_agent_timestamp.', 'Serverless: 0-50,000 RU/s.\nNo capacity planning needed.'],
        ['Dashboard', 'Container Apps auto-scales 0-N replicas.\nWebSocket sticky sessions.', 'Plotly WebGL: 2000-point rolling\nbuffers, client-side rendering.'],
    ]
    add_table(slide, rows, MARGIN, cy + Inches(0.25), lw)

    # Provisioning checklist
    py = cy + Inches(1.9)
    section_box(slide, 'Provisioning Checklist ✅', [
        '✅ requirements.txt — all dependencies pinned',
        '✅ Dockerfile — containerized, multi-stage build',
        '✅ python main.py --smoke-test — 5-tick pipeline validation',
        '✅ USE_AZURE=true — single flag activates full Azure stack',
        '✅ Container Registry + Container Apps deployed & running',
        '✅ Cosmos DB + AI Search provisioned with data loaded',
        '✅ Application Insights configured with custom metrics',
    ], MARGIN + lw + Inches(0.15), py - Inches(0.5), lw, Inches(2.1), title_bg=GREEN, bg=LIGHT_GREEN)
    footer(slide, 17)


def slide_18_mcp_reuse(prs):
    """MCP Packaging & Cross-Industry Transfer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'MCP Packaging & Cross-Industry Transfer',
                   'Model Context Protocol servers + domain-agnostic safety pattern')

    # MCP table
    my = y + Inches(0.08)
    tf = txbox(slide, MARGIN, my, CW, Inches(0.25))
    para(tf, '5 MCP Servers — Typed Input/Output Schemas, Composable with Any Orchestrator:', sz=11, color=NAVY, bold=True)
    rows = [
        ['MCP Server', 'Tools Exposed', 'Input Schema', 'Output Schema', 'Orchestrator Compatibility'],
        ['sensor-mcp', 'check_thresholds(state)\nget_alerts()', 'ReactorState\n(11 fields)', 'Alert[]\n(dimension, severity, value)',
         'AI Foundry, Semantic Kernel,\nAutoGen, LangChain'],
        ['risk-mcp', 'assess_risk(state, alerts)\nget_score_history()', 'State + Alert[]', 'RiskScore (0-100) +\ntrend + reasoning',
         'AI Foundry, AutoGen,\nCopilot Studio'],
        ['decision-mcp', 'decide(state, risk_score)\nget_action_log()', 'State + RiskScore', 'Decision[] with source\nattribution (rule/llm)',
         'Copilot Studio, Foundry,\nSemantic Kernel'],
        ['evacuation-mcp', 'plan_evacuation(pop, routes)\nget_progress()', 'Population +\nRoute configs', 'EvacPlan + ETAs +\nbus assignments',
         'Semantic Kernel,\nAutoGen'],
        ['comms-mcp', 'draft_alert(state, decisions)\nmodel_plume(params)', 'State + Decision[]\n+ wind params', 'AlertMsg +\nPlumeData[] (5-300km)',
         'AI Foundry, AutoGen,\nLangChain'],
    ]
    add_table(slide, rows, MARGIN, my + Inches(0.25), CW)

    # Guardrail-Union as reusable pattern
    gy = my + Inches(2.15)
    code_box(slide,
        'THE REUSABLE PATTERN — Guardrail-Union Architecture (applies to ANY safety-critical domain):\n'
        '  1. Telemetry → Sensor Agent (deterministic threshold monitoring)\n'
        '  2. Alerts → Risk Agent (rule + AI blend scoring)\n'
        '  3. Risk Score → Decision Agent:\n'
        '     ├─ Rules run FIRST (safety floor — deterministic, non-negotiable)\n'
        '     ├─ AI runs SECOND (proactive advisory — adds early warnings)\n'
        '     └─ Final = UNION(rules, AI) — AI can ADD actions, never REMOVE\n'
        '  4. Hard guardrails are non-negotiable (no override path exists)\n'
        '  5. Graceful degradation to pure deterministic rules on any AI failure\n'
        '  6. Adversarial agent stress-tests all decisions under realistic human pressure',
        MARGIN, gy, CW, Inches(1.65))

    # Cross-industry adaptation
    ay = gy + Inches(1.8)
    tf2 = txbox(slide, MARGIN, ay, CW, Inches(0.25))
    para(tf2, 'Cross-Industry Adaptation — 6 Modules to Change, Everything Else Transfers:', sz=11, color=NAVY, bold=True)

    rows2 = [
        ['Industry', 'ReactorState →', 'SCRAM →', 'Radiation →', 'Dyatlov →', 'IAEA →'],
        ['Oil & Gas', 'WellState (PSI, flow)', 'Emergency Well Shutdown', 'H₂S / blowout', 'Production Manager pressure', 'BSEE + Coast Guard'],
        ['Aviation', 'FlightState (alt, speed)', 'Go-Around / Diversion', 'Turbulence / bird strike', 'Captain authority gradient', 'ATC + NTSB'],
        ['Pharmaceuticals', 'BatchState (pH, temp)', 'Batch Rejection / Lockdown', 'Contamination risk', 'Manufacturing quota', 'FDA + recall notice'],
        ['Power Grid', 'GridState (V, Hz, MW)', 'Load Shedding / Islanding', 'Frequency deviation', 'Dispatcher congestion', 'NERC + ISO report'],
    ]
    add_table(slide, rows2, MARGIN, ay + Inches(0.25), CW)
    footer(slide, 18)


def slide_19_validation(prs):
    """Real-World Validation & KPIs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = header_bar(slide, 'Real-World Validation & Impact',
                   'Same architecture pattern deployed in production for 45 million people')

    # Validation table
    vy = y + Inches(0.08)
    tf = txbox(slide, MARGIN, vy, CW, Inches(0.25))
    para(tf, 'Industry Deployments Using the Same Guardrail-Union Pattern:', sz=11, color=NAVY, bold=True)
    rows = [
        ['Deployment', 'Organization', 'Scale', 'PRIPYAT-1986 Equivalent', 'Status'],
        ['MISO Azure AI\nGrid Platform', 'Microsoft + MISO\n(Jan 2026)', '45 million people\n15 US states',
         'Full pipeline: Sensor → Risk → Decision.\nReal-time grid resilience operations.', 'Production\n(live)'],
        ['eGridGPT', 'NREL\n(US National Lab)', 'National grid\nresearch scale',
         'RiskAgent + DecisionAgent pattern.\nAI-augmented grid analysis.', 'Research\n(active)'],
        ['ACWA Power\nPredictive Maint.', 'Microsoft +\nACWA Power', 'Multiple power\nplants',
         'SensorAgent + RiskAgent pattern.\nPredictive maintenance scoring.', 'Production\n(deployed)'],
        ['Control Room\nof the Future', 'Capgemini\n(CIGRE 2025)', 'Framework\nstandard',
         'Complete CRoF architecture framework.\nOur reference architecture.', 'Published\n(CIGRE)'],
    ]
    add_table(slide, rows, MARGIN, vy + Inches(0.25), CW)

    # Our KPI Results
    ky = vy + Inches(1.9)
    tf2 = txbox(slide, MARGIN, ky, CW, Inches(0.25))
    para(tf2, 'PRIPYAT-1986 Simulation Results — Measurable KPIs:', sz=12, color=NAVY, bold=True)

    kpi_box(slide, '38 min', 'Earlier Emergency Shutdown\nvs. actual explosion at 01:23:40', MARGIN, ky + Inches(0.3))
    kpi_box(slide, '35 hrs', 'Earlier Civilian Evacuation\nvs. Soviet order Apr 27 14:00', Inches(3.5), ky + Inches(0.3))
    kpi_box(slide, '100%', 'Safety Protocol Compliance\nvs. 0% on April 26, 1986', Inches(6.5), ky + Inches(0.3))
    kpi_box(slide, '<$0.10', 'Cost Per Simulation Run\non Azure OpenAI GPT-4o', Inches(9.5), ky + Inches(0.3))

    # Second row of KPIs
    kpi_box(slide, '5 ticks', 'Smoke Test Validation\nFull pipeline in <10 seconds', MARGIN, ky + Inches(1.3), val_color=ACCENT)
    kpi_box(slide, '<5 min', 'Deployment Time\nDocker → Container Apps', Inches(3.5), ky + Inches(1.3), val_color=ACCENT)
    kpi_box(slide, '~$5-15', 'Total Azure Spend\nof $350 hackathon budget', Inches(6.5), ky + Inches(1.3), val_color=ACCENT)
    kpi_box(slide, '6 agents', '3 LLM + 3 Deterministic\n100% functional without API', Inches(9.5), ky + Inches(1.3), val_color=ACCENT)

    # Innovation callout
    iy = ky + Inches(2.3)
    section_box(slide, 'What Makes This Different from a Generic Chatbot', [
        'Multi-dimensional compound risk detection (invisible to single-threshold systems)',
        'Guardrail-Union: rules + AI merged, not OR\'d — LLM adds value, never removes safety',
        'Adversarial stress testing built into the architecture (DyatlovAgent)',
        'Dual timeline engine: historical replay + AI counterfactual in parallel',
        'MCP-packaged agents composable with any Microsoft orchestrator (Foundry, Semantic Kernel, AutoGen)',
        'Same pattern validated in production by MISO + Microsoft for 45M people — not theoretical',
    ], MARGIN, iy, CW, Inches(1.8), title_bg=NAVY)
    footer(slide, 19)


def slide_20_pitch(prs):
    """The Pitch — Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SW, SH, NAVY)

    # Teal stripe
    rect(slide, Inches(0), Inches(2.4), SW, Pt(3), TEAL)

    # Title
    tf = txbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.8))
    para(tf, 'PRIPYAT-1986', sz=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    para(tf, 'The Architecture is the Product', sz=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Big quote
    qbox = rrect(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(2.0), DARK_BLUE)
    qtf = qbox.text_frame; qtf.word_wrap = True
    qtf.margin_left = Inches(0.4); qtf.margin_top = Inches(0.15); qtf.margin_right = Inches(0.4)
    para(qtf, '', sz=6, color=WHITE)
    para(qtf, '"PRIPYAT-1986 is not a nuclear simulator.', sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    para(qtf, 'It\'s a reusable safety-critical AI architecture', sz=16, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    para(qtf, 'that uses Chernobyl as its most dramatic training dataset.', sz=14, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
    para(qtf, '', sz=6, color=WHITE)
    para(qtf, 'Connect it to your SCADA API → it runs on your grid tomorrow.', sz=13, color=WHITE, align=PP_ALIGN.CENTER)
    para(qtf, 'Connect it to your BSEE feed → it monitors your wells.', sz=13, color=WHITE, align=PP_ALIGN.CENTER)
    para(qtf, 'The architecture is the product. The disaster is just the proof."', sz=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Summary boxes
    bw = Inches(3.5); bh = Inches(0.75); by = Inches(5.0)
    items = [
        ('Track A: GitHub + Azure', '7 Azure services deployed'),
        ('6 Agents, 3 LLM-powered', 'Guardrail-Union architecture'),
        ('< $0.10 per run', '100% rule-based fallback'),
    ]
    for i, (title, sub) in enumerate(items):
        x = Inches(0.7) + i * (bw + Inches(0.3))
        box = rrect(slide, x, by, bw, bh, TEAL)
        btf = box.text_frame; btf.word_wrap = True
        btf.margin_left = Inches(0.1); btf.margin_top = Inches(0.05)
        para(btf, title, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        para(btf, sub, sz=10, color=RGBColor(0xDD, 0xFF, 0xDD), align=PP_ALIGN.CENTER)

    # Team
    tf3 = txbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.4))
    para(tf3, 'Team WattAgents: Shruti Kate • Amar Shikhar • V. Sreenivas • Deendayalan Balaji • M. Navenaa • Shyamdinesh',
         sz=11, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    tf4 = txbox(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.3))
    para(tf4, 'github.com/amarshikhar/pripyat-1986', sz=10, color=MGRAY, align=PP_ALIGN.CENTER)

    # Footer
    tf5 = txbox(slide, Inches(1), Inches(6.9), Inches(11), Inches(0.3))
    para(tf5, 'Company Confidential © Capgemini 2026. All rights reserved.', sz=7, color=MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_01_title(prs)
    slide_02_exec_summary(prs)
    slide_03_not_a_simulator(prs)
    slide_04_use_case(prs)
    slide_05_timeline(prs)
    slide_06_pipeline(prs)
    slide_07_guardrail_union(prs)
    slide_08_risk_scoring(prs)
    slide_09_physics(prs)
    slide_10_dyatlov(prs)
    slide_11_azure_arch(prs)
    slide_12_foundry_copilot(prs)
    slide_13_security(prs)
    slide_14_responsible_ai(prs)
    slide_15_operator_ux(prs)
    slide_16_reliability_cost(prs)
    slide_17_scalability(prs)
    slide_18_mcp_reuse(prs)
    slide_19_validation(prs)
    slide_20_pitch(prs)

    out = os.path.join(os.path.dirname(__file__), 'docs', 'WATTAGENTS_PRIPYAT_1986_v3.pptx')
    prs.save(out)
    print(f'Saved: {out}')
    print(f'Slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
