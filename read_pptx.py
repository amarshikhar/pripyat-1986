from pptx import Presentation

prs = Presentation(r'C:\Users\shamar\pripyat-1986\docs\WATTAGENTS- utitly sector (2).pptx')
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} (layout: {slide.slide_layout.name}) ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'  {text}')
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                joined = ' | '.join(cells)
                print(f'  | {joined} |')
    print()
